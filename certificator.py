import shutil
import os
import comtypes.client
from pptx import Presentation
import pandas as pd
import re
import zipfile

# Constants
COLLEGE_NAME = "SENA"
TEMPLATE_FILENAME = 'templateSENA.pptx'
LISTNAMES_EXCEL_FILE = 'data.xlsx'
PROCESSED_FILES_FOLDER = f"./{COLLEGE_NAME}/"
FILENAME = f"CERTIFICATE_{COLLEGE_NAME}_"
OLDTEXT_PLACEHOLDER = '_X_'
PDF_FORMAT_CONSTANT = 32
ZIP_FILENAME = f"{COLLEGE_NAME}_Certificates.zip"

def generate_pptx_files():
    os.mkdir(PROCESSED_FILES_FOLDER)
    prs = Presentation(TEMPLATE_FILENAME)
    new_files = []

    data = pd.read_excel(LISTNAMES_EXCEL_FILE, header=None, usecols=[1], skiprows=1)
    new_texts = data.iloc[:, 0].tolist()

    for name in new_texts:
        name = str(name)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if OLDTEXT_PLACEHOLDER in run.text:
                                run.text = run.text.replace(OLDTEXT_PLACEHOLDER, name)
        name = re.sub(r'\W+', '_', name)
        new_file = PROCESSED_FILES_FOLDER + FILENAME + name + '.pptx'
        prs.save(new_file)
        new_files.append(new_file)
        prs = Presentation(TEMPLATE_FILENAME)
    
    return new_files

def convert_to_pdf(pptx_files):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = True
    pdf_files = []

    for new_file in pptx_files:
        ppt_path = os.path.abspath(new_file)
        pdf_path = os.path.abspath(new_file[:-5] + '.pdf')
        deck = powerpoint.Presentations.Open(ppt_path)
        deck.SaveAs(pdf_path, PDF_FORMAT_CONSTANT)
        deck.Close()
        os.remove(ppt_path)
        pdf_files.append(pdf_path)
    
    powerpoint.Quit()
    return pdf_files

def create_zip(pdf_files):
    with zipfile.ZipFile(ZIP_FILENAME, mode='w') as zip_file:
        for pdf_file in pdf_files:
            # Calculate relative path with respect to PROCESSED_FILES_FOLDER
            relative_path = os.path.relpath(pdf_file, PROCESSED_FILES_FOLDER)
            zip_file.write(pdf_file, arcname=relative_path)
    
    shutil.rmtree(PROCESSED_FILES_FOLDER)

def main():
    print(f"Generating certificates using {TEMPLATE_FILENAME} template and names included in {LISTNAMES_EXCEL_FILE} file.")
    pptx_files = generate_pptx_files()
    print(f"Successfully created pptx temporary files. Generating PDF files.")
    pdf_files = convert_to_pdf(pptx_files)
    print(f"Successfully created pdf files. Generating ZIP file.")
    create_zip(pdf_files)
    print(f"Successfully created {ZIP_FILENAME}")

if __name__ == "__main__":
    main()
